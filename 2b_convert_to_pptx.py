import cv2
from lxml import etree
import os
from pptx import Presentation
from pptx.util import Inches

from utils import *

THUMBNAIL_FILENAME = f"{DIRECTORY}/thumbnail.png"
VIDEO_FILENAME = f"{DIRECTORY}/video.mov"

videos = read_output_videos()
height, width, _ = videos[0][0].shape

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

for idx, frames in enumerate(videos):
    cv2.imwrite(THUMBNAIL_FILENAME, frames[0])
    video_writer = cv2.VideoWriter(VIDEO_FILENAME, cv2.VideoWriter_fourcc(*"mp4v"), 60, (width, height))
    for frame in frames:
        video_writer.write(frame)
    video_writer.release()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    movie = slide.shapes.add_movie(VIDEO_FILENAME, 0, 0, prs.slide_width, prs.slide_height, poster_frame_image=THUMBNAIL_FILENAME)
    tree = movie._element.getparent().getparent().getnext().getnext()
    timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == "cond"][0]
    timing.set("delay", "0")

    print(f"\033[30;1mProcessed video #{idx + 1}\033[0m")

prs.save(f"{DIRECTORY}/output.pptx")
os.unlink(VIDEO_FILENAME)
os.unlink(THUMBNAIL_FILENAME)

print(f"\033[32;1mDone!\033[0m")
