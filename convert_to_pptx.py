import cv2
from lxml import etree
import os
from pptx import Presentation
from pptx.util import Inches

DIRECTORY = os.path.realpath(os.path.dirname(__file__))
OUTPUT_DIRECTORY = f"{DIRECTORY}/output"
THUMBNAIL_FILENAME = f"{OUTPUT_DIRECTORY}/thumbnail.png"

cap_filenames = ["output/" + filename for filename in sorted(os.listdir(OUTPUT_DIRECTORY)) if filename.endswith(".mp4")]

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

for cap_filename in cap_filenames:
    cap = cv2.VideoCapture(cap_filename)
    _, frame = cap.read()
    cap.release()
    cv2.imwrite(THUMBNAIL_FILENAME, frame)

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    movie = slide.shapes.add_movie(cap_filename, 0, 0, prs.slide_width, prs.slide_height, poster_frame_image=THUMBNAIL_FILENAME)
    tree = movie._element.getparent().getparent().getnext().getnext()
    timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == "cond"][0]
    timing.set("delay", "0")

    print(f"\033[30;1mLoaded {cap_filename}\033[0m")

prs.save(f"{DIRECTORY}/output.pptx")
os.unlink(THUMBNAIL_FILENAME)

print(f"\033[32;1mDone!\033[0m")
